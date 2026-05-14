"""Unit tests for D-10: OfficeExtractor.

Strategy: Docling's ``DocumentConverter`` with the default pipeline
(no VLM yet — P1 territory) → ``result.document.export_to_markdown()``.
Handles ``.docx`` / ``.pptx`` / ``.xlsx``.

Test fixtures build minimal artefacts via python-docx / python-pptx /
openpyxl (all transitive deps of Docling) so the test file is self-
contained inside ``tmp_path``.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from corpus_forge.extractors import ExtractedDocument, Extractor
from corpus_forge.extractors.office import OfficeExtractor

# ── Fixture helpers ──────────────────────────────────────────────────


def _build_docx(tmp_path: Path, name: str = "report.docx") -> Path:
    from docx import Document

    d = Document()
    d.add_heading("Executive Summary", level=1)
    d.add_paragraph("This is the body of the first paragraph in the report.")
    d.add_paragraph("Second paragraph contains follow-up content.")
    target = tmp_path / name
    d.save(str(target))
    return target


def _build_pptx(tmp_path: Path, name: str = "deck.pptx") -> Path:
    from pptx import Presentation

    pres = Presentation()
    layout = pres.slide_layouts[5]  # title-only
    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Title Slide"
    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Second Slide"
    target = tmp_path / name
    pres.save(str(target))
    return target


def _build_xlsx(tmp_path: Path, name: str = "sheet.xlsx") -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "alpha"
    ws["B2"] = 1
    ws["A3"] = "beta"
    ws["B3"] = 2
    target = tmp_path / name
    wb.save(str(target))
    return target


# ── Tests ────────────────────────────────────────────────────────────


def test_extractor_protocol_conformance():
    ex: Extractor = OfficeExtractor()
    assert isinstance(ex.supported_extensions, tuple)


def test_supported_extensions():
    ex = OfficeExtractor()
    assert set(ex.supported_extensions) == {".docx", ".pptx", ".xlsx"}


def test_extract_docx(tmp_path: Path):
    p = _build_docx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert isinstance(doc, ExtractedDocument)
    assert doc.chunker_hint == "markdown"
    assert "Executive Summary" in doc.text
    assert "first paragraph" in doc.text


def test_extract_docx_metadata(tmp_path: Path):
    p = _build_docx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.metadata.get("extractor") == "docling"
    assert doc.metadata.get("format") == "docx"


def test_extract_docx_labels(tmp_path: Path):
    p = _build_docx(tmp_path)
    doc = OfficeExtractor().extract(p)
    labels = {(ns, val) for ns, val in doc.labels}
    assert ("format", "docx") in labels
    assert ("extractor", "docling") in labels


def test_extract_pptx(tmp_path: Path):
    p = _build_pptx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.chunker_hint == "markdown"
    assert "Title Slide" in doc.text
    assert "Second Slide" in doc.text


def test_extract_pptx_metadata_format(tmp_path: Path):
    p = _build_pptx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.metadata.get("format") == "pptx"


def test_extract_pptx_metadata_page_count(tmp_path: Path):
    """Docling reports slide count via ``num_pages()`` on the document."""
    p = _build_pptx(tmp_path)
    doc = OfficeExtractor().extract(p)
    # Docling's PPTX path reports the slide count.
    assert doc.metadata.get("page_count") == 2


def test_extract_xlsx(tmp_path: Path):
    p = _build_xlsx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.chunker_hint == "markdown"
    # Docling exports tables as Markdown tables.
    assert "Name" in doc.text
    assert "alpha" in doc.text
    assert "beta" in doc.text


def test_extract_xlsx_metadata_format(tmp_path: Path):
    p = _build_xlsx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.metadata.get("format") == "xlsx"


def test_extract_xlsx_labels(tmp_path: Path):
    p = _build_xlsx(tmp_path)
    doc = OfficeExtractor().extract(p)
    labels = {(ns, val) for ns, val in doc.labels}
    assert ("format", "xlsx") in labels


def test_extract_language_is_none(tmp_path: Path):
    p = _build_docx(tmp_path)
    doc = OfficeExtractor().extract(p)
    assert doc.language is None


def test_lazy_import_does_not_load_docling_on_module_import():
    import subprocess
    import sys

    script = (
        "import sys; "
        "import corpus_forge.extractors.office as m; "
        "assert 'docling' not in sys.modules, sorted(sys.modules); "
        "print('ok')"
    )
    result = subprocess.run(
        [sys.executable, "-c", script], capture_output=True, text=True, check=False
    )
    assert result.returncode == 0, result.stderr
    assert result.stdout.strip() == "ok"


@pytest.mark.parametrize("ext", [".docx", ".pptx", ".xlsx"])
def test_registry_wires_office_extractor(tmp_path: Path, ext: str):
    from corpus_forge.extractors import register_default_extractors

    reg = register_default_extractors(config=None)
    p = tmp_path / f"x{ext}"
    p.write_bytes(b"placeholder")
    extractor = reg.get_for(p)
    assert isinstance(extractor, OfficeExtractor)
