"""Deterministic synthetic fixture-corpus generator — Phase D / D-17.

Builds the in-repo fixture tree under
``tests/fixtures/multi_format_corpus/`` that the multi-format e2e
integration test (``tests/integration/test_multi_format_ingest_e2e.py``)
runs against. The tree is committed to the repo as-is; this script
exists so contributors can regenerate it from a clean checkout.

Determinism contract
--------------------
* No timestamps, UUIDs, or per-run randomness anywhere in file content.
  Every output file is byte-identical across runs and machines.
* Office formats (DOCX / PPTX / XLSX) and EPUB are zip containers; we
  pin every embedded ``CT_/`` timestamp to ``2000-01-01 00:00:00`` and
  every zip-entry mtime to the same epoch.
* Re-running the script is idempotent — existing files are overwritten
  with the same bytes; the tree never accumulates stale entries.

Usage
-----
::

    uv sync --group dev --extra multi-format --extra code
    uv run python scripts/build_fixture_corpus.py

The script lives in ``scripts/`` (not in ``corpus_forge/``) because its
heavy build-time deps (reportlab, python-docx, python-pptx, openpyxl,
Pillow, ebooklib) live in the dev dependency group only. Production
ingest never imports this module.
"""

from __future__ import annotations

import io
import os
import shutil
import zipfile
from datetime import datetime, timezone
from pathlib import Path

# ── Deterministic clock ─────────────────────────────────────────────────────
# Anything that writes a timestamp into a file we control routes through
# this constant. Pinned to 2000-01-01 00:00:00 UTC so the repo doesn't
# accumulate epoch-2026 timestamps that confuse archaeological greps.
_EPOCH = datetime(2000, 1, 1, tzinfo=timezone.utc)
_EPOCH_TUPLE = (2000, 1, 1, 0, 0, 0)

# Force zip-tool writers to use _EPOCH_TUPLE rather than time.time() so the
# .docx / .pptx / .xlsx / .epub bytes stay reproducible. python-docx and its
# kin call ``zipfile.ZipFile.writestr`` under the hood; that path picks up
# ``time.localtime()`` unless ZIP_DEFLATED is paired with an explicit
# ZipInfo.date_time. Easiest catch-all: monkey-patch zipfile's ``time``
# import to a callable that returns _EPOCH_TUPLE.

_orig_zipinfo_init = zipfile.ZipInfo.__init__


def _deterministic_zipinfo_init(  # type: ignore[no-untyped-def]
    self,
    filename: str = "NoName",
    date_time: tuple[int, int, int, int, int, int] = _EPOCH_TUPLE,
) -> None:
    """Force ZipInfo.date_time to the fixture epoch unless caller overrides."""
    _orig_zipinfo_init(self, filename, date_time)


zipfile.ZipInfo.__init__ = _deterministic_zipinfo_init  # type: ignore[assignment]

# ── Filesystem helpers ──────────────────────────────────────────────────────


_ROOT = Path(__file__).resolve().parent.parent / "tests" / "fixtures" / "multi_format_corpus"


def _write_bytes(rel_path: str, payload: bytes) -> Path:
    """Write ``payload`` to ``rel_path`` under the fixture root, mkdir-p parents.

    Idempotent: overwrites any existing file. Returns the absolute path
    written.
    """
    target = _ROOT / rel_path
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_bytes(payload)
    # Pin mtime so `git diff` doesn't see metadata-only changes between
    # script runs. Use the fixture epoch.
    os.utime(target, (_EPOCH.timestamp(), _EPOCH.timestamp()))
    return target


def _write_text(rel_path: str, text: str) -> Path:
    """Write UTF-8 text to ``rel_path`` (newline-normalised to LF)."""
    # Normalise line endings so the bytes are stable on Windows checkouts
    # too. The integration test reads UTF-8; the LF normalisation makes
    # content_hash predictable.
    normalised = text.replace("\r\n", "\n").replace("\r", "\n")
    if not normalised.endswith("\n"):
        normalised += "\n"
    return _write_bytes(rel_path, normalised.encode("utf-8"))


# ── Prose family (markdown / txt / rst / org / tex / adoc) ─────────────────


def build_prose() -> None:
    _write_text(
        "prose/intro.md",
        """# Introduction

Lorem ipsum dolor sit amet, consectetur adipiscing elit.
Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.

## Section 1

Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.

## Section 2

Duis aute irure dolor in reprehenderit in voluptate velit esse cillum dolore.
""",
    )
    _write_text(
        "prose/notes.txt",
        """Plain-text notes for the fixture corpus.
Two short paragraphs separated by a blank line.

A second paragraph confirming the PlainTextExtractor wires the
chunker_hint = "passthrough" path.
""",
    )
    _write_text(
        "prose/frontmatter.md",
        """---
title: Frontmatter Example
tags: [fixture, multi-format]
---

# Frontmatter Example

Markdown body following YAML frontmatter. The PassthroughMarkdownExtractor
emits this verbatim; ChunkerDispatcher routes it through MarkdownChunker.
""",
    )
    _write_text(
        "prose/tex-snippet.tex",
        r"""\documentclass{article}
\begin{document}
\section{Tiny LaTeX Fixture}

A trivial LaTeX document so the PlainTextExtractor's ``.tex`` extension
maps to chunker_hint = "passthrough". Equation: $E = mc^2$.
\end{document}
""",
    )


# ── PDF family (digital-only; OCR is P1) ───────────────────────────────────


def build_pdfs() -> None:
    from reportlab.lib.pagesizes import LETTER  # noqa: PLC0415
    from reportlab.lib.units import inch  # noqa: PLC0415
    from reportlab.pdfgen import canvas  # noqa: PLC0415

    # ── digital-single-col.pdf ──────────────────────────────────────────
    buf1 = io.BytesIO()
    c = canvas.Canvas(buf1, pagesize=LETTER)
    # ``setDateFormatter`` clears reportlab's xref timestamp output.
    c.setCreator("corpus-forge fixture builder")
    c.setTitle("Single-column digital PDF fixture")
    c.setAuthor("corpus-forge")
    c.setSubject("D-17 fixture")
    text = c.beginText(1 * inch, 10 * inch)
    text.setFont("Helvetica", 11)
    for line in [
        "Single-column digital PDF fixture.",
        "",
        "This PDF exercises the PdfDigitalExtractor's pymupdf4llm path.",
        "It is intentionally tiny and contains only ASCII text in one column,",
        "ensuring the extractor returns non-empty markdown without escalating",
        "to the (Wave 5) VLM OCR pipeline.",
        "",
        "Section one body. Lorem ipsum dolor sit amet.",
        "Section two body. Consectetur adipiscing elit.",
    ]:
        text.textLine(line)
    c.drawText(text)
    c.showPage()
    c.save()
    pdf1 = _strip_pdf_timestamps(buf1.getvalue())
    _write_bytes("pdf/digital-single-col.pdf", pdf1)

    # ── digital-two-col-equations.pdf ───────────────────────────────────
    # The plan explicitly says: real multi-column rendering is not needed;
    # plain-text equation symbols are sufficient for the test's "non-empty
    # extraction" assertion. We keep a single column with math glyphs.
    buf2 = io.BytesIO()
    c = canvas.Canvas(buf2, pagesize=LETTER)
    c.setCreator("corpus-forge fixture builder")
    c.setTitle("Two-column equation PDF fixture")
    c.setAuthor("corpus-forge")
    c.setSubject("D-17 fixture")
    text = c.beginText(1 * inch, 10 * inch)
    text.setFont("Helvetica", 11)
    for line in [
        "Two-column-equations PDF fixture.",
        "",
        "sum_{i=0}^{n} i = n(n+1)/2",
        "integral_{0}^{1} x dx = 1/2",
        "E = m c^2",
        "",
        "This PDF carries math-like ASCII so the extractor's text-layer pass",
        "yields non-empty markdown without invoking the VLM escalation path.",
    ]:
        text.textLine(line)
    c.drawText(text)
    c.showPage()
    c.save()
    pdf2 = _strip_pdf_timestamps(buf2.getvalue())
    _write_bytes("pdf/digital-two-col-equations.pdf", pdf2)


def _strip_pdf_timestamps(pdf_bytes: bytes) -> bytes:
    """Strip reportlab's non-deterministic per-run metadata.

    reportlab embeds three pieces of wall-clock-bound state in every PDF:

    * ``/CreationDate (D:20260514153000-04'00')`` — current time.
    * ``/ModDate (D:...)`` — current time.
    * ``/ID [<32-hex><32-hex>]`` — per-build random document id.

    For byte-stable fixtures we rewrite all three to fixed values. The
    substitutions are byte-level on tightly-bounded regexes.
    """
    import re  # noqa: PLC0415

    out = re.compile(rb"/(CreationDate|ModDate) \(D:[^)]*\)").sub(
        rb"/\1 (D:20000101000000+00'00')", pdf_bytes
    )
    # ``/ID`` field is two hex strings inside angle brackets, e.g.
    # ``/ID \n[<deadbeef...><deadbeef...>]``. Pin to a fixed pair so
    # repeat builds emit identical bytes.
    fixed_id = b"00000000000000000000000000000000"
    out = re.compile(rb"/ID\s*\[<[0-9a-fA-F]{32}><[0-9a-fA-F]{32}>\]").sub(
        b"/ID \n[<" + fixed_id + b"><" + fixed_id + b">]", out
    )
    return out


# ── HTML family (readability happy path + nav/ads stripping) ────────────


def build_html() -> None:
    _write_text(
        "html/simple-article.html",
        """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>Simple article</title>
</head>
<body>
  <article>
    <h1>Simple article fixture</h1>
    <p>A short article body with a single paragraph.</p>
    <p>readability-lxml should extract this as the main content with no
    boilerplate stripping required.</p>
    <h2>Subsection</h2>
    <p>Another paragraph so the markdownify pass produces at least two
    block-level elements.</p>
  </article>
</body>
</html>
""",
    )
    _write_text(
        "html/nav-and-ads.html",
        """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>Article with nav and ads</title>
</head>
<body>
  <header>
    <nav>
      <ul>
        <li><a href="/home">Home</a></li>
        <li><a href="/about">About</a></li>
        <li><a href="/contact">Contact</a></li>
      </ul>
    </nav>
  </header>
  <aside class="ads">
    <p>SUBSCRIBE NOW! Limited-time offer!</p>
    <p>Buy our newsletter today!</p>
  </aside>
  <main>
    <article>
      <h1>Article with surrounding noise</h1>
      <p>readability-lxml is responsible for stripping the navigation,
      sidebar ads, and footer boilerplate, leaving only this article
      body.</p>
      <p>The article needs at least two paragraphs and a heading so the
      extractor produces a non-trivial main-content block. Otherwise
      readability falls back to including the whole page.</p>
      <p>Filler third paragraph to give the readability scorer enough
      density to confidently pick this article element. Lorem ipsum
      dolor sit amet consectetur adipiscing elit.</p>
    </article>
  </main>
  <footer>
    <p>Copyright 2026 corpus-forge</p>
  </footer>
</body>
</html>
""",
    )


# ── EPUB (ebooklib happy path) ──────────────────────────────────────────


def build_epub() -> None:
    # ebooklib calls zipfile.ZipFile under the hood; our monkey-patched
    # ZipInfo.__init__ pins the embedded mtimes to the fixture epoch.
    from ebooklib import epub  # noqa: PLC0415

    book = epub.EpubBook()
    book.set_identifier("corpus-forge-fixture-small-book")
    book.set_title("Small fixture book")
    book.set_language("en")
    book.add_author("corpus-forge")

    chapters = []
    for i, (title, body) in enumerate(
        [
            (
                "Chapter 1 — Genesis",
                "<h1>Chapter 1 — Genesis</h1>"
                "<p>The first chapter of the synthetic EPUB. Two short "
                "paragraphs so the EpubExtractor produces a multi-section "
                "markdown body.</p>"
                "<p>A second paragraph in the same chapter.</p>",
            ),
            (
                "Chapter 2 — Conflict",
                "<h1>Chapter 2 — Conflict</h1>"
                "<p>Second chapter; same shape as the first. The chapters "
                "are joined by the extractor with a horizontal rule "
                "separator (<code>---</code>).</p>"
                "<p>A second paragraph to ensure chunker has enough text.</p>",
            ),
        ],
        start=1,
    ):
        ch = epub.EpubHtml(title=title, file_name=f"chap_{i}.xhtml", lang="en")
        ch.content = body
        book.add_item(ch)
        chapters.append(ch)

    book.toc = tuple(chapters)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", *chapters]

    out = io.BytesIO()
    epub.write_epub(out, book, {})
    # ebooklib writes a fresh dcterms:modified timestamp into content.opf
    # on every build — pipe through _zip_normalise to scrub it.
    _write_bytes("epub/small-book.epub", _zip_normalise(out.getvalue()))


# ── Office (Docling extracts; we author with python-docx/-pptx, openpyxl) ─


def build_office() -> None:
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_heading("Report fixture", level=1)
    doc.add_paragraph(
        "A trivial DOCX file generated by python-docx for the multi-format "
        "fixture corpus. Docling extracts this as markdown."
    )
    doc.add_heading("Findings", level=2)
    doc.add_paragraph(
        "Lorem ipsum dolor sit amet. The OfficeExtractor's Docling pipeline "
        "should turn this into one or two markdown blocks."
    )
    docx_buf = io.BytesIO()
    doc.save(docx_buf)
    _write_bytes("office/report.docx", _zip_normalise(docx_buf.getvalue()))

    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415

    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slides fixture"
    slide.placeholders[1].text = "Synthetic PPTX for multi-format corpus"
    # Content slide
    content_layout = prs.slide_layouts[1]
    s2 = prs.slides.add_slide(content_layout)
    s2.shapes.title.text = "Findings"
    body_tf = s2.placeholders[1].text_frame
    body_tf.text = "First bullet point"
    body_tf.add_paragraph().text = "Second bullet point"
    pptx_buf = io.BytesIO()
    prs.save(pptx_buf)
    _write_bytes("office/slides.pptx", _zip_normalise(pptx_buf.getvalue()))

    from openpyxl import Workbook  # noqa: PLC0415

    wb = Workbook()
    ws = wb.active
    ws.title = "tiny-sheet"
    ws.append(["name", "value", "note"])
    ws.append(["alpha", 1, "first row"])
    ws.append(["beta", 2, "second row"])
    ws.append(["gamma", 3, "third row"])
    xlsx_buf = io.BytesIO()
    wb.save(xlsx_buf)
    _write_bytes("office/tiny-sheet.xlsx", _zip_normalise(xlsx_buf.getvalue()))


def _zip_normalise(zip_bytes: bytes) -> bytes:
    """Repack a zip archive with sorted entries, pinned mtime, and scrubbed
    XML timestamps inside ``core.xml`` / ``content.opf``-style payloads.

    Office files (.docx / .pptx / .xlsx) and EPUBs are zip containers
    whose entries embed wall-clock timestamps in metadata XML. We can't
    rely on per-library setters alone (openpyxl, python-docx, ebooklib
    all phrase the field differently), so we scrub at the bytes layer
    using tightly-bounded regex substitutions:

    * ``<dcterms:created ...>2026-05-15T02:49:17Z</dcterms:created>``
    * ``<dcterms:modified ...>...</dcterms:modified>``
    * ``<meta property="dcterms:modified">...</meta>`` (EPUB content.opf)
    """
    import re  # noqa: PLC0415

    fixed = "2000-01-01T00:00:00Z"
    dcterms_open_close = re.compile(
        rb"(<dcterms:(?:created|modified)\b[^>]*>)\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}Z(</dcterms:\w+>)"
    )
    epub_meta_modified = re.compile(
        rb'(<meta\s+property="dcterms:modified"[^>]*>)'
        rb"\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}Z"
        rb"(</meta>)"
    )

    def _scrub(data: bytes) -> bytes:
        data = dcterms_open_close.sub(rb"\g<1>" + fixed.encode() + rb"\g<2>", data)
        data = epub_meta_modified.sub(rb"\g<1>" + fixed.encode() + rb"\g<2>", data)
        return data

    in_buf = io.BytesIO(zip_bytes)
    out_buf = io.BytesIO()
    with zipfile.ZipFile(in_buf, "r") as zin, zipfile.ZipFile(
        out_buf, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for name in sorted(zin.namelist()):
            data = zin.read(name)
            # Only XML/OPF payloads need timestamp scrubbing — images,
            # binary metadata, etc. stay untouched.
            if name.endswith((".xml", ".opf")):
                data = _scrub(data)
            info = zipfile.ZipInfo(filename=name, date_time=_EPOCH_TUPLE)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o644 << 16
            zout.writestr(info, data)
    return out_buf.getvalue()


# ── Notebook ────────────────────────────────────────────────────────────


def build_notebook() -> None:
    import json  # noqa: PLC0415

    # Stable cell IDs so the nbformat 5.1+ "MissingIDFieldWarning" stays
    # silent and the bytes stay reproducible across runs (no UUIDs).
    nb = {
        "cells": [
            {
                "id": "fixture-cell-001",
                "cell_type": "markdown",
                "metadata": {},
                "source": [
                    "# Notebook fixture\n",
                    "\n",
                    "A trivial notebook with two markdown and two code cells.\n",
                ],
            },
            {
                "id": "fixture-cell-002",
                "cell_type": "code",
                "execution_count": None,
                "metadata": {},
                "outputs": [],
                "source": ["def hello() -> str:\n", "    return 'hello'\n"],
            },
            {
                "id": "fixture-cell-003",
                "cell_type": "markdown",
                "metadata": {},
                "source": ["## Section\n", "\n", "Second markdown cell.\n"],
            },
            {
                "id": "fixture-cell-004",
                "cell_type": "code",
                "execution_count": None,
                "metadata": {},
                "outputs": [],
                "source": ["print(hello())\n"],
            },
        ],
        "metadata": {
            "kernelspec": {
                "display_name": "Python 3",
                "language": "python",
                "name": "python3",
            },
            "language_info": {
                "name": "python",
                "version": "3.13.0",
            },
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    text = json.dumps(nb, indent=1, ensure_ascii=False, sort_keys=True)
    _write_text("notebook/analysis.ipynb", text)


# ── Data family (csv / toml / json / srt) ──────────────────────────────


def build_data() -> None:
    _write_text(
        "data/records.csv",
        """id,name,score
1,alpha,0.91
2,beta,0.42
3,gamma,0.77
4,delta,0.13
""",
    )
    _write_text(
        "data/config.toml",
        """[server]
host = "127.0.0.1"
port = 8080

[features]
multi_format = true
ocr = false
""",
    )
    _write_text(
        "data/manifest.json",
        """{
  "name": "fixture-corpus",
  "version": "1.0.0",
  "description": "Synthetic fixture set for the multi-format e2e test",
  "files": ["records.csv", "config.toml"]
}
""",
    )
    _write_text(
        "data/transcript.srt",
        """1
00:00:00,000 --> 00:00:02,500
Hello and welcome to the fixture corpus.

2
00:00:02,500 --> 00:00:05,000
This subtitle file exercises the SubtitleExtractor.
""",
    )


# ── Code family (one file per language, syntactically valid trivia) ────


def build_code() -> None:
    # Python
    _write_text(
        "code/python/module.py",
        '''"""Trivial Python module fixture."""


def hello() -> str:
    """Return a greeting."""
    return "hello, world"


class Greeter:
    """A trivial class."""

    def greet(self, name: str) -> str:
        return f"hello, {name}"
''',
    )
    _write_text(
        "code/python/package/__init__.py",
        '''"""Trivial Python package fixture."""

from .util import compute  # noqa: F401
''',
    )
    _write_text(
        "code/python/package/util.py",
        '''"""Trivial util module."""


def compute(x: int, y: int) -> int:
    """Add two integers."""
    return x + y
''',
    )

    # C / C++
    _write_text(
        "code/cpp/main.cpp",
        """#include <iostream>
#include "util.h"

int main() {
    std::cout << add(1, 2) << std::endl;
    return 0;
}
""",
    )
    _write_text(
        "code/cpp/util.h",
        """#pragma once

int add(int a, int b);
""",
    )
    _write_text(
        "code/cpp/util.cpp",
        """#include "util.h"

int add(int a, int b) {
    return a + b;
}
""",
    )
    _write_text(
        "code/c/lib.c",
        """#include "lib.h"

int square(int x) {
    return x * x;
}
""",
    )
    _write_text(
        "code/c/lib.h",
        """#ifndef LIB_H
#define LIB_H

int square(int x);

#endif
""",
    )

    # JS / TS / TSX
    _write_text(
        "code/js-ts/app.ts",
        """export function greet(name: string): string {
    return `hello, ${name}`;
}
""",
    )
    _write_text(
        "code/js-ts/react.tsx",
        """import * as React from "react";

export function Hello(): React.ReactElement {
    return <div>hello, world</div>;
}
""",
    )
    _write_text(
        "code/js-ts/server.js",
        """function main() {
    console.log("hello, world");
}

main();
""",
    )

    # Go
    _write_text(
        "code/go/main.go",
        """package main

import "fmt"

func main() {
    fmt.Println("hello, world")
}
""",
    )
    _write_text(
        "code/go/internal/handler.go",
        """package internal

func Handle(input string) string {
    return "handled: " + input
}
""",
    )

    # Rust
    _write_text(
        "code/rust/main.rs",
        """fn main() {
    println!("hello, world");
}
""",
    )
    _write_text(
        "code/rust/lib.rs",
        """pub fn add(a: i32, b: i32) -> i32 {
    a + b
}
""",
    )

    # Java
    _write_text(
        "code/java/App.java",
        """public class App {
    public static void main(String[] args) {
        System.out.println("hello, world");
    }
}
""",
    )

    # Kotlin / Scala
    _write_text(
        "code/kotlin-scala/App.kt",
        """fun main() {
    println("hello, world")
}
""",
    )
    _write_text(
        "code/kotlin-scala/App.scala",
        """object App {
  def main(args: Array[String]): Unit = {
    println("hello, world")
  }
}
""",
    )

    # Ruby
    _write_text(
        "code/ruby/app.rb",
        """def hello
  "hello, world"
end

puts hello
""",
    )

    # BEAM (Elixir / Erlang)
    _write_text(
        "code/beam/hello.ex",
        """defmodule Hello do
  def world do
    "hello, world"
  end
end
""",
    )
    _write_text(
        "code/beam/mod.erl",
        """-module(mod).
-export([hello/0]).

hello() ->
    "hello, world".
""",
    )
    _write_text(
        "code/beam/mod.hrl",
        """-define(GREETING, "hello, world").
""",
    )

    # Prolog (exercises tree-sitter long-tail / fallback path)
    _write_text(
        "code/prolog/rules.pl",
        """parent(tom, bob).
parent(bob, alice).

ancestor(X, Y) :- parent(X, Y).
ancestor(X, Y) :- parent(X, Z), ancestor(Z, Y).
""",
    )

    # Haskell / OCaml
    _write_text(
        "code/haskell-ocaml/Main.hs",
        """module Main where

main :: IO ()
main = putStrLn "hello, world"
""",
    )
    _write_text(
        "code/haskell-ocaml/main.ml",
        """let () = print_endline "hello, world"
""",
    )

    # Lisps / Clojure
    _write_text(
        "code/lisp-clj/core.clj",
        """(ns fixture.core)

(defn hello [] "hello, world")

(println (hello))
""",
    )
    _write_text(
        "code/lisp-clj/demo.lisp",
        """(defun hello () "hello, world")

(format t "~a~%" (hello))
""",
    )
    _write_text(
        "code/lisp-clj/demo.scm",
        """(define (hello) "hello, world")

(display (hello))
(newline)
""",
    )

    # Shells
    _write_text(
        "code/shell/install.sh",
        """#!/bin/sh
echo "installing"
""",
    )
    _write_text(
        "code/shell/deploy.bash",
        """#!/bin/bash
set -e
echo "deploying"
""",
    )
    _write_text(
        "code/shell/fish.fish",
        """function hello
    echo "hello, world"
end

hello
""",
    )

    # Web (CSS / HTML / SQL)
    _write_text(
        "code/web/styles.css",
        """body {
    font-family: monospace;
    color: #333;
}

.heading {
    font-weight: bold;
}
""",
    )
    _write_text(
        "code/web/page.html",
        """<!DOCTYPE html>
<html>
  <head><title>Page</title></head>
  <body><p>hello, world</p></body>
</html>
""",
    )
    _write_text(
        "code/web/query.sql",
        """SELECT id, name
FROM users
WHERE active = 1
ORDER BY name;
""",
    )

    # Exotic (zig / nim / crystal / r / julia / swift / dart / nix)
    _write_text(
        "code/exotic/tiny.zig",
        """const std = @import("std");

pub fn main() void {
    std.debug.print("hello, world\\n", .{});
}
""",
    )
    _write_text(
        "code/exotic/mod.nim",
        """proc hello(): string =
  "hello, world"

echo hello()
""",
    )
    _write_text(
        "code/exotic/app.cr",
        """def hello
  "hello, world"
end

puts hello
""",
    )
    _write_text(
        "code/exotic/plot.r",
        """hello <- function() {
  return("hello, world")
}

print(hello())
""",
    )
    _write_text(
        "code/exotic/run.jl",
        """function hello()
    return "hello, world"
end

println(hello())
""",
    )
    _write_text(
        "code/exotic/app.swift",
        """func hello() -> String {
    return "hello, world"
}

print(hello())
""",
    )
    _write_text(
        "code/exotic/ui.dart",
        """String hello() {
  return 'hello, world';
}

void main() {
  print(hello());
}
""",
    )
    _write_text(
        "code/exotic/default.nix",
        """{ pkgs ? import <nixpkgs> {} }:

pkgs.stdenv.mkDerivation {
  name = "fixture";
  src = ./.;
}
""",
    )

    # Build / extension-less files (Makefile, Dockerfile, dotfiles)
    _write_text(
        "code/build/Makefile",
        """.PHONY: hello

hello:
\t@echo "hello, world"
""",
    )
    _write_text(
        "code/build/Dockerfile",
        """FROM alpine:3.20

RUN echo "hello, world" > /etc/greeting

CMD ["cat", "/etc/greeting"]
""",
    )
    _write_text(
        "code/build/.gitignore",
        """*.pyc
__pycache__/
.venv/
""",
    )
    _write_text(
        "code/build/.editorconfig",
        """root = true

[*]
indent_style = space
indent_size = 4
end_of_line = lf
""",
    )


# ── README ─────────────────────────────────────────────────────────────


def build_readme() -> None:
    _write_text(
        "README.md",
        """# `tests/fixtures/multi_format_corpus/`

Synthetic fixture tree for the multi-format end-to-end integration test
(`tests/integration/test_multi_format_ingest_e2e.py`).

## How to regenerate

The tree is **checked in** as committed bytes. Only regenerate when
adding new file families or fixing a determinism bug:

```bash
uv sync --group dev --extra multi-format --extra code
uv run python scripts/build_fixture_corpus.py
git diff tests/fixtures/multi_format_corpus
```

A clean `git diff` after a regen run is the determinism contract.

## Layout

```
multi_format_corpus/
├── prose/      Markdown / plain-text / LaTeX (PassthroughMarkdownExtractor + PlainTextExtractor)
├── pdf/        Digital PDFs (PdfDigitalExtractor)
├── html/       Articles with and without nav/ads noise (HtmlExtractor)
├── epub/       Tiny EPUB (EpubExtractor)
├── office/     DOCX / PPTX / XLSX (OfficeExtractor — Docling)
├── notebook/   Jupyter notebook (NotebookExtractor)
├── data/       CSV / TOML / JSON / SRT (CsvExtractor / StructuredDataExtractor / SubtitleExtractor)
└── code/       Source files for ~20 languages + Makefile / Dockerfile / dotfiles (CodeExtractor)
```

## Notes

* No timestamps, UUIDs, or per-run randomness in any file content.
* Office and EPUB files are repacked zip containers with sorted entries
  and pinned epoch mtimes, so the bytes are stable across machines.
* PDFs have their `/CreationDate` and `/ModDate` stripped to the same
  fixture epoch (2000-01-01).
* A local `.gitignore` lives at this fixture root with negation rules
  re-including `code/build/` — the project-root `.gitignore` excludes
  `build/` everywhere, which would otherwise drop the Makefile /
  Dockerfile / `.gitignore` / `.editorconfig` fixtures.
* Wave 6 / P1 fixtures (`images/` and `pdf/scanned-paper.pdf`) are
  generated by `build_p1_ocr_fixtures()` in the same script. They feed
  the live-Ollama and Mistral OCR end-to-end tests under
  `tests/integration/test_ocr_local_e2e.py` and
  `tests/integration/test_ocr_remote_e2e.py`. Rendering is pinned to
  Pillow's bundled DejaVu Sans Bold so the bytes stay reproducible.
""",
    )


def _clean_p1_only_subtrees() -> None:
    """Remove subtrees that this script does not own (P1 only).

    Wave 6 / P1 now owns ``images/`` and ``pdf/scanned-paper.pdf`` and
    generates them via :func:`build_p1_ocr_fixtures` below. Kept as a
    no-op for backwards compatibility (and to document that the P0
    builder does NOT touch P1 subtrees unilaterally).
    """


# ── P1 / Wave 6 fixtures (OCR escalation tests) ─────────────────────────


def _pillow_default_font(*, size: int):  # type: ignore[no-untyped-def]
    """Return Pillow's bundled DejaVu Sans Bold at the requested size.

    Pillow ships this TTF as in-memory bytes (BytesIO), so the rendering
    output is byte-identical across machines + Pillow ≥10.x. Avoids
    pulling in system fonts that vary per OS / per user.
    """
    from PIL import ImageFont  # noqa: PLC0415

    return ImageFont.load_default(size=size)


def _render_text_image(
    *,
    size: tuple[int, int],
    lines: list[tuple[int, int, str, int]],
    background: tuple[int, int, int] = (255, 255, 255),
    text_color: tuple[int, int, int] = (0, 0, 0),
):  # type: ignore[no-untyped-def]
    """Render ``lines`` of ``(x, y, text, font_size)`` onto a white image.

    Returns a :class:`PIL.Image.Image` in RGB mode. Determinism is
    guaranteed by:

    - Pillow's bundled DejaVu Sans Bold font (BytesIO, no system
      fallback);
    - integer-pixel coordinates (no antialiasing seeded by sub-pixel
      offsets);
    - a fixed RGB background colour.
    """
    from PIL import Image, ImageDraw  # noqa: PLC0415

    img = Image.new("RGB", size, background)
    draw = ImageDraw.Draw(img)
    for x, y, text, font_size in lines:
        font = _pillow_default_font(size=font_size)
        draw.text((x, y), text, fill=text_color, font=font)
    return img


def _render_diagram_image():  # type: ignore[no-untyped-def]
    """Two-circle "A → B" diagram on a white 400x400 canvas."""
    from PIL import Image, ImageDraw  # noqa: PLC0415

    size = (400, 400)
    img = Image.new("RGB", size, (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Connecting line first so the circles overdraw the ends.
    draw.line([(110, 200), (290, 200)], fill=(0, 0, 0), width=3)

    # Left circle "A"
    draw.ellipse([(40, 140), (160, 260)], outline=(0, 0, 0), width=3)
    font = _pillow_default_font(size=40)
    draw.text((90, 175), "A", fill=(0, 0, 0), font=font)

    # Right circle "B"
    draw.ellipse([(240, 140), (360, 260)], outline=(0, 0, 0), width=3)
    draw.text((290, 175), "B", fill=(0, 0, 0), font=font)

    return img


def build_p1_images() -> None:
    """Generate the three P1 image fixtures (PNG / JPEG / WebP).

    Each writer pins encoder options that affect byte output:

    - PNG: ``optimize=True``, ``compress_level=9`` for stable zlib stream.
    - JPEG: ``quality=85``, ``optimize=True``, ``progressive=False``.
    - WebP: ``lossless=True``, ``quality=90``, ``method=6``.

    Combined with the deterministic rendering helpers above, two
    consecutive ``python scripts/build_fixture_corpus.py`` runs must
    produce byte-identical files (``git diff`` empty).
    """
    # ── screenshot.png — "code editor" mock with readable monospace text ──
    screenshot = _render_text_image(
        size=(600, 400),
        lines=[
            (20, 20, "def hello(name: str) -> str:", 16),
            (20, 50, "    return f'hello, {name}'", 16),
            (20, 90, "if __name__ == '__main__':", 16),
            (20, 120, "    print(hello('world'))", 16),
            (20, 180, "# Expected output:", 16),
            (20, 210, "# hello, world", 16),
            (20, 280, "Editor screenshot fixture for OCR.", 14),
            (20, 310, "The VLM should transcribe this text.", 14),
        ],
    )
    screenshot_path = _ROOT / "images" / "screenshot.png"
    screenshot_path.parent.mkdir(parents=True, exist_ok=True)
    screenshot.save(
        screenshot_path,
        format="PNG",
        optimize=True,
        compress_level=9,
    )
    os.utime(screenshot_path, (_EPOCH.timestamp(), _EPOCH.timestamp()))

    # ── photo-of-receipt.jpg — synthetic receipt with line items + total ─
    receipt = _render_text_image(
        size=(400, 600),
        lines=[
            (60, 30, "CORPUS MARKET", 24),
            (90, 70, "123 Fixture Lane", 14),
            (110, 95, "Test City, AA", 14),
            (40, 150, "Item            Qty   Price", 14),
            (40, 175, "---------------------------", 14),
            (40, 200, "Coffee           1    3.50", 14),
            (40, 225, "Bagel            2    4.00", 14),
            (40, 250, "Muffin           1    2.84", 14),
            (40, 275, "Juice            1    2.00", 14),
            (40, 320, "Subtotal             12.34", 14),
            (40, 345, "Tax                   0.00", 14),
            (40, 380, "TOTAL                12.34", 18),
            (40, 440, "Thank you for shopping!", 14),
        ],
    )
    receipt_path = _ROOT / "images" / "photo-of-receipt.jpg"
    receipt.save(
        receipt_path,
        format="JPEG",
        quality=85,
        optimize=True,
        progressive=False,
    )
    os.utime(receipt_path, (_EPOCH.timestamp(), _EPOCH.timestamp()))

    # ── diagram.webp — minimal "A → B" diagram ──────────────────────────
    diagram = _render_diagram_image()
    diagram_path = _ROOT / "images" / "diagram.webp"
    diagram.save(
        diagram_path,
        format="WebP",
        lossless=True,
        quality=90,
        method=6,
    )
    os.utime(diagram_path, (_EPOCH.timestamp(), _EPOCH.timestamp()))


def build_p1_scanned_pdf() -> None:
    """Generate a 1-page image-only PDF (no text layer) for OCR escalation.

    Construction:

    1. Render a "scanned paper" looking image via Pillow (white BG, body
       text drawn as pixels).
    2. Embed the rasterised image as the entire page content of a 1-page
       PDF via reportlab's ``Canvas.drawInlineImage``.
    3. Strip reportlab's ``/CreationDate`` and ``/ModDate`` to the
       fixture epoch.

    Because the page content is a raster image (not a text run), the
    embedded text layer is empty — ``pymupdf4llm.helpers.pymupdf_rag.to_markdown``
    returns near-zero characters, the ``PdfDigitalExtractor`` flags
    ``sparse_text_layer=True``, and Tier 2 escalation fires.
    """
    from reportlab.lib.pagesizes import LETTER  # noqa: PLC0415
    from reportlab.lib.utils import ImageReader  # noqa: PLC0415
    from reportlab.pdfgen import canvas  # noqa: PLC0415

    # Generate the "scan" — keep the canvas modest (612x792 at 1 px per pt
    # is fine; the VLM only needs to read the text, not produce a true
    # 300 DPI scan).
    page_image = _render_text_image(
        size=(612, 792),
        lines=[
            (60, 60, "Scanned Paper Fixture", 24),
            (60, 110, "Phase D / Wave 6 OCR escalation test.", 14),
            (60, 140, "This PDF has no embedded text layer.", 14),
            (60, 170, "Every glyph above this line is rendered", 14),
            (60, 195, "as image pixels.", 14),
            (60, 250, "Section 1 — Introduction", 18),
            (60, 290, "Lorem ipsum dolor sit amet, consectetur", 14),
            (60, 315, "adipiscing elit. The PdfDigitalExtractor", 14),
            (60, 340, "should detect the empty text layer and", 14),
            (60, 365, "escalate to the active VLM backend.", 14),
            (60, 420, "Section 2 — Expected Behaviour", 18),
            (60, 460, "qwen2.5vl:7b transcribes this fixture to", 14),
            (60, 485, "non-empty Markdown via /api/generate.", 14),
            (60, 540, "End of scanned-paper fixture.", 14),
        ],
    )

    # Encode the Pillow image to in-memory PNG bytes so reportlab can
    # embed it deterministically (no temp files, no per-run filenames).
    png_buf = io.BytesIO()
    page_image.save(png_buf, format="PNG", optimize=True, compress_level=9)
    png_buf.seek(0)

    pdf_buf = io.BytesIO()
    c = canvas.Canvas(pdf_buf, pagesize=LETTER)
    c.setCreator("corpus-forge fixture builder")
    c.setTitle("Scanned-paper PDF fixture")
    c.setAuthor("corpus-forge")
    c.setSubject("E-07 fixture")
    # Letter is 612x792 pts; draw the image at the full page size.
    c.drawImage(ImageReader(png_buf), 0, 0, width=612, height=792)
    c.showPage()
    c.save()
    pdf_bytes = _strip_pdf_timestamps(pdf_buf.getvalue())
    _write_bytes("pdf/scanned-paper.pdf", pdf_bytes)


def build_p1_ocr_fixtures() -> None:
    """Generate every P1 / Wave 6 OCR-related fixture.

    Idempotent. Overwrites the three image files and the scanned-paper
    PDF with byte-identical output across consecutive runs.
    """
    build_p1_scanned_pdf()
    build_p1_images()


def main() -> None:
    # Idempotency: leave _ROOT in place but overwrite known children.
    # We don't ``shutil.rmtree(_ROOT)`` because the user may have local
    # P1 fixtures under ``images/`` that aren't this script's business.
    _ROOT.mkdir(parents=True, exist_ok=True)

    build_readme()
    build_prose()
    build_pdfs()
    build_html()
    build_epub()
    build_office()
    build_notebook()
    build_data()
    build_code()
    build_p1_ocr_fixtures()
    _clean_p1_only_subtrees()


if __name__ == "__main__":
    main()
    # No print — keep the script silent under `uv run`.
